import asyncio
import edge_tts
import os
import json 
from pptx import Presentation
from openai import OpenAI
import tempfile
import shutil
import time
from tqdm import tqdm
from PIL import Image
import numpy as np
import win32com.client as win32
from moviepy.video.VideoClip import ImageClip
from moviepy.video.io.ImageSequenceClip import ImageSequenceClip
from moviepy import concatenate_videoclips
from moviepy import AudioFileClip

async def generate_tts_async(text, path, voice="zh-CN-XiaoxiaoNeural", rate="+0%"):
    communicate = edge_tts.Communicate(text=text, voice=voice, rate=rate)
    await communicate.save(path)

def generate_voice_for_slides(ppt_content, output_dir="voice_audio"):
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    loop = asyncio.get_event_loop()
    
    # 生成封面的语音文件
    cover_text = ppt_content['title'] + "\n" + "汉赋自习室"  # 假设封面副标题是 "汉赋自习室"
    cover_filename = os.path.join(output_dir, f"slide_000.mp3")
    print(f"正在为封面合成语音...")
    loop.run_until_complete(generate_tts_async(cover_text, cover_filename))
    loop = asyncio.get_event_loop()
    for idx, page in enumerate(ppt_content['pages']):
        # 正确使用字符串拼接（移除多余的反斜杠转义）
        text = page['title'] + "\n" + "\n".join(
            f"{blk['title']}。{blk['description']}" for blk in page['content']
        )
        filename = os.path.join(output_dir, f"slide_{idx+1:03d}.mp3")
        print(f"正在为第 {idx+2} 页合成语音...")
        loop.run_until_complete(generate_tts_async(text, filename))
    return output_dir
    
def estimate_ppt_pages_from_content(text, target_chars_per_page=200, min_pages=5, max_pages=15):
    char_count = len(text.replace('\\n', '').replace(' ', ''))
    estimated = max(min_pages, min(max_pages, char_count // target_chars_per_page + 1))
    return estimated

def return_llm(query, history=[], user_stop_words=[]): 
    client = OpenAI(api_key="sk-jrvcstpnxgsoxwsaehhwmasgggywdvgeeaurlltclfckgfcn", 
    base_url="https://api.siliconflow.cn/v1")

    messages = [{'role': 'system', 'content': 'You are a helpful assistant.'}]
    for hist in history:
        messages.append({'role': 'user', 'content': hist[0]})
        messages.append({'role': 'assistant', 'content': hist[1]})
    messages.append({'role': 'user', 'content': query})

    response = client.chat.completions.create(
        model="Qwen/Qwen2.5-72B-Instruct",
        messages=messages,
        stream=True
    )

    full_response = ""
    for chunk in response:
        if chunk.choices[0].delta.content:
            full_response += chunk.choices[0].delta.content
    return full_response

def generate_ppt_content(topic, pages):
    output_format = json.dumps({
        "title": "example title",
        "pages": [
            {
                "title": "title for page 1",
                "content": [
                    {
                        "title": "title for paragraph 1",
                        "description": "detail for paragraph 1",
                    }
                ],
            }
        ],
    }, ensure_ascii=True)
    
    prompt = f'''我要准备1个关于{topic}的PPT，要求一共写{pages}页，请你根据主题生成详细内容，每一点写三到五句话就行。
    严格按这个JSON格式输出{output_format}，只能返回JSON，不要有任何其他字符，请认真检查生成的JSON文本的合法性，如引号未闭合等问题。且JSON不要用```包裹，内容要用中文。'''

    prompt_topic = f'''请根据我给出的主题：{topic},重新生成一个书面化的，
                       可以作为课件ppt标题的名称，直接给出该名称，不要有任何其他字符与标点符号。
                       示例：
                       问题：凯恩斯主义是啥啊
                       错误回答：凯恩斯主义概述与解析
                       正确回答：凯恩斯主义
                       只能返回JSON，且JSON不要用```包裹，内容要用中文。'''

    topic_refined_json = json.loads(return_llm(prompt_topic))
    topic_refined = topic_refined_json if isinstance(topic_refined_json, str) else list(topic_refined_json.values())[0]

    ppt_content = json.loads(return_llm(prompt))
    return ppt_content, topic_refined


def generate_ppt_file(topic, ppt_content):
    ppt = Presentation()
    
    slide = ppt.slides.add_slide(ppt.slide_layouts[0]) 
    slide.placeholders[0].text = ppt_content['title']
    slide.placeholders[1].text = "汉赋自习室"
    
    print('总共%d页...' % len(ppt_content['pages']))
    for i, page in enumerate(ppt_content['pages']):
        print('生成第%d页:%s' % (i + 1, page['title']))
        slide = ppt.slides.add_slide(ppt.slide_layouts[1]) 
        slide.placeholders[0].text = page['title']
        for sub_content in page['content']:
            sub_title = slide.placeholders[1].text_frame.add_paragraph()
            sub_title.text, sub_title.level = sub_content['title'], 1
            sub_description = slide.placeholders[1].text_frame.add_paragraph()
            sub_description.text, sub_description.level = sub_content['description'], 2
    
    ppt_path = f'{topic}.pptx'
    ppt.save(ppt_path)
    return os.path.abspath(ppt_path)

def emu_to_pixels(emu, dpi=96):
    return int((emu / 914400) * dpi)

def add_background_to_ppt(ppt_path_background, bg_image_path):
    prs = Presentation(ppt_path_background)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    slide_width_px = emu_to_pixels(slide_width)
    slide_height_px = emu_to_pixels(slide_height)

    bg_img = Image.open(bg_image_path).convert("RGB")
    bg_img = bg_img.resize((slide_width_px, slide_height_px))
    temp_img_path = "resized_bg_temp.jpg"
    bg_img.save(temp_img_path)

    for slide in prs.slides:
        bg_shape = slide.shapes.add_picture(temp_img_path, 0, 0, width=slide_width, height=slide_height)
        slide.shapes._spTree.remove(bg_shape._element)
        slide.shapes._spTree.insert(2, bg_shape._element)

    output_path_background = f'{topic}_background.pptx'
    prs.save(output_path_background)
    os.remove(temp_img_path)
    return os.path.abspath(output_path_background)

def get_slide_durations_from_content(ppt_content, seconds_per_char=0.15, min_duration=1):
    durations = []
    # 封面时长
    cover_text = ppt_content['title'] + "\n" + "汉赋自习室"
    cover_char_count = len(cover_text.strip())
    cover_duration = max(min_duration, cover_char_count * seconds_per_char)
    durations.append(cover_duration)
    for page in ppt_content['pages']:
        total_text = page['title']
        for block in page['content']:
            total_text += block['title'] + block['description']
        char_count = len(total_text.strip())
        duration = max(min_duration, char_count * seconds_per_char)
        durations.append(duration)
    return durations

def ppt_to_video(ppt_path, output_path, durations=None, fps=30, transition_duration=1):
    if not os.path.exists(ppt_path):
        print(f"错误: 文件 '{ppt_path}' 不存在。")
        return
    
    temp_dir = tempfile.mkdtemp()
    image_dir = os.path.join(temp_dir, "images")
    os.makedirs(image_dir, exist_ok=True)
    
    print(f"正在导出PPT幻灯片为图像...")
    try:
        powerpoint = win32.gencache.EnsureDispatch('PowerPoint.Application')
        presentation = powerpoint.Presentations.Open(ppt_path)
        slides = presentation.Slides
        
        for i, slide in enumerate(tqdm(slides, desc="导出幻灯片")):
            slide_img_path = os.path.join(image_dir, f"slide_{i:03d}.png")
            slide.Export(slide_img_path, 'PNG')
        
        presentation.Close()
        powerpoint.Quit()
    except Exception as e:
        print(f"打开PPT文件时出错: {e}")
        shutil.rmtree(temp_dir, ignore_errors=True)
        return
    
    print("正在生成视频片段...")
    clips = []
    image_files = sorted(os.listdir(image_dir))
    num_slides = len(image_files)
    audio_dir = "voice_audio"  # 路径需与语音生成处保持一致

    for i in range(num_slides):
        current_img_path = os.path.join(image_dir, image_files[i])
        slide_duration = durations[i] if durations and i < len(durations) else 5
        current_clip = ImageClip(current_img_path).with_duration(slide_duration)
        clips.append(current_clip)
        
        if i < num_slides - 1:
            next_img_path = os.path.join(image_dir, image_files[i+1])
            transition_clip = create_transition_clip(
                current_img_path, next_img_path,
                duration=transition_duration, fps=fps
            )
        current_clip = ImageClip(current_img_path).with_duration(slide_duration)

        audio_path = os.path.join(audio_dir, f"slide_{i:03d}.mp3")
        if os.path.exists(audio_path):
            audio_clip = AudioFileClip(audio_path)
            slide_duration = audio_clip.duration  # 使用音频实际长度
            current_clip = current_clip.with_duration(slide_duration).with_audio(audio_clip)
        clips.append(current_clip)
    
    
    print("正在合并视频片段...")
    final_clip = concatenate_videoclips(clips, method="compose")
    
    print("正在保存视频...")
    final_clip.write_videofile(output_path, fps=fps, codec="libx264")
    
    time.sleep(1)
    shutil.rmtree(temp_dir, ignore_errors=True)
    
    print(f"🎉 视频已成功保存至: {output_path}")

def create_transition_clip(img_path1, img_path2, duration, fps):
    img1 = Image.open(img_path1).convert('RGB')
    img2 = Image.open(img_path2).convert('RGB')
    
    width, height = max(img1.size, img2.size)
    img1 = img1.resize((width, height), Image.LANCZOS)
    img2 = img2.resize((width, height), Image.LANCZOS)
    
    frames = []
    num_frames = int(duration * fps)
    
    for i in range(num_frames):
        alpha = i / (num_frames - 1)
        blended = Image.blend(img1, img2, alpha)
        frame_array = np.array(blended)
        frames.append(frame_array)
    
    return ImageSequenceClip(frames, fps=fps)

if __name__ == '__main__':
    topic = input('输入主题:')
    initial_pages = 11
    
    ppt_content, refined_topic = generate_ppt_content(topic, initial_pages)

# 统计文字总数
    all_text = ppt_content['title'] + ''.join(
        page['title'] + ''.join(block['title'] + block['description'] for block in page['content'])
        for page in ppt_content['pages']
    )

# 根据内容重新估算页数
    estimated_pages = estimate_ppt_pages_from_content(all_text)

    # 若实际生成页数和预期不一致，重新生成
    if len(ppt_content['pages']) != estimated_pages:
        print(f"🔁 内容偏少或偏多，重新生成 {estimated_pages} 页内容...")
        ppt_content, refined_topic = generate_ppt_content(topic, estimated_pages)

        ppt_path = generate_ppt_file(refined_topic, ppt_content)
    
        ppt_path_background = add_background_to_ppt(
            ppt_path_background=ppt_path,
            bg_image_path="R-C.jpg",
        )

    durations = get_slide_durations_from_content(ppt_content, seconds_per_char=0.15, min_duration=1)
    durations[0] = 5
    generate_voice_for_slides(ppt_content)
    output_path = os.path.splitext(ppt_path)[0] + ".mp4"
    ppt_to_video(
        ppt_path_background,
        output_path,
        durations=durations,
        fps=30,
        transition_duration=1
    )

    os.remove(ppt_path)
    os.remove(ppt_path_background)
    print(f"临时生成的PPT文件 {ppt_path} 已删除。")
